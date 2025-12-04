"""
Unit tests for PPTXReportGenerator.

Tests:
- Basic PPTX generation
- Data type conversion (DataFrame, dict, list)
- Chart generation (bar, column, line, pie)
- Multi-slide pagination
- Custom theme colors
- Font customization
- Configuration validation
- Error handling

Status: Phase 3B - PPTX Support (1M-360 FINAL PHASE)

Code Reuse: Test patterns from test_excel_generator.py (100% consistency)
"""

from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation

from extract_transform_platform.reports import (
    PPTXReportConfig,
    PPTXReportGenerator,
    ReportConfig,
)


class TestPPTXReportGenerator:
    """Test suite for PPTXReportGenerator."""

    @pytest.fixture
    def sample_data(self) -> pd.DataFrame:
        """Sample DataFrame for testing."""
        return pd.DataFrame(
            {
                "Region": ["North", "South", "East", "West"],
                "Sales": [45000, 52000, 38000, 61000],
                "Target": [50000, 50000, 50000, 50000],
                "Growth": ["10%", "15%", "5%", "20%"],
            }
        )

    @pytest.fixture
    def output_path(self, tmp_path: Path) -> Path:
        """Temporary output path for PPTX files."""
        return tmp_path / "test_report.pptx"

    def test_generate_basic_pptx(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test basic PPTX generation with default config."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="Sales Report Q4")

        result_path = generator.generate(sample_data, output_path, config)

        # Verify file exists
        assert result_path.exists()
        assert result_path == output_path.resolve()

        # Verify file is valid PPTX
        prs = Presentation(str(result_path))
        assert len(prs.slides) >= 2  # Title + at least 1 data slide

    def test_pptx_with_chart(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test PPTX generation with chart slide."""
        config = PPTXReportConfig(
            title="Sales Report with Chart", include_chart=True, chart_type="bar"
        )
        generator = PPTXReportGenerator()
        result_path = generator.generate(sample_data, output_path, config)

        assert result_path.exists()

        # Verify presentation has title + data + chart slides
        prs = Presentation(str(result_path))
        assert len(prs.slides) >= 3  # Title + data + chart

    def test_pptx_multi_slide_pagination(self, output_path: Path) -> None:
        """Test pagination with large dataset."""
        # Create dataset with 25 rows (should split into 3 slides with max_rows=10)
        large_data = pd.DataFrame(
            {"Item": [f"Item {i}" for i in range(25)], "Value": range(25)}
        )

        config = PPTXReportConfig(title="Large Dataset Report", max_rows_per_slide=10)
        generator = PPTXReportGenerator()
        result_path = generator.generate(large_data, output_path, config)

        assert result_path.exists()

        # Verify slide count: Title + 3 data slides
        prs = Presentation(str(result_path))
        assert len(prs.slides) == 4

    def test_pptx_chart_types(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test different chart types (bar, column, line, pie)."""
        for chart_type in ["bar", "column", "line", "pie"]:
            config = PPTXReportConfig(
                title=f"{chart_type.capitalize()} Chart Report",
                include_chart=True,
                chart_type=chart_type,
            )
            generator = PPTXReportGenerator()
            chart_output = output_path.parent / f"test_{chart_type}.pptx"
            result_path = generator.generate(sample_data, chart_output, config)

            assert result_path.exists()

            # Verify presentation has chart slide
            prs = Presentation(str(result_path))
            assert len(prs.slides) >= 3  # Title + data + chart

    def test_pptx_custom_theme_color(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test custom theme color."""
        config = PPTXReportConfig(
            title="Custom Color Report", theme_color="#FF5733"  # Orange-red
        )
        generator = PPTXReportGenerator()
        result_path = generator.generate(sample_data, output_path, config)

        assert result_path.exists()

    def test_pptx_custom_font(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test custom font settings."""
        config = PPTXReportConfig(
            title="Custom Font Report", font_name="Arial", font_size=16
        )
        generator = PPTXReportGenerator()
        result_path = generator.generate(sample_data, output_path, config)

        assert result_path.exists()

    def test_get_supported_features(self) -> None:
        """Test supported features list."""
        generator = PPTXReportGenerator()
        features = generator.get_supported_features()

        assert "slides" in features
        assert "tables" in features
        assert "charts" in features
        assert "themes" in features
        assert "layouts" in features
        assert "multi_slide" in features
        assert "colors" in features

    def test_invalid_output_extension_raises(
        self, sample_data: pd.DataFrame, tmp_path: Path
    ) -> None:
        """Test invalid extension raises ValueError."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="Test")
        invalid_path = tmp_path / "test.txt"

        with pytest.raises(ValueError, match="must have .pptx extension"):
            generator.generate(sample_data, invalid_path, config)

    def test_wrong_config_type_raises(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test wrong config type raises TypeError."""
        generator = PPTXReportGenerator()
        wrong_config = ReportConfig(title="Test")

        with pytest.raises(TypeError, match="requires PPTXReportConfig"):
            generator.generate(sample_data, output_path, wrong_config)

    def test_empty_data_raises(self, output_path: Path) -> None:
        """Test empty data raises ValueError."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="Test")

        with pytest.raises(ValueError, match="cannot be empty"):
            generator.generate(pd.DataFrame(), output_path, config)

    def test_none_data_raises(self, output_path: Path) -> None:
        """Test None data raises ValueError."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="Test")

        with pytest.raises(ValueError, match="cannot be None"):
            generator.generate(None, output_path, config)

    def test_unsupported_data_type_raises(self, output_path: Path) -> None:
        """Test unsupported data type raises TypeError."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="Test")

        with pytest.raises(TypeError, match="Unsupported data type"):
            generator.generate(12345, output_path, config)

    def test_dict_to_dataframe(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test dict to DataFrame conversion."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="Dict Test")

        # Single row dict
        data_dict = {"Region": "North", "Sales": 45000}
        result_path = generator.generate(data_dict, output_path, config)

        assert result_path.exists()

    def test_list_of_dicts_to_dataframe(self, output_path: Path) -> None:
        """Test list of dicts to DataFrame conversion."""
        generator = PPTXReportGenerator()
        config = PPTXReportConfig(title="List Test")

        # List of dicts
        data_list = [
            {"Region": "North", "Sales": 45000},
            {"Region": "South", "Sales": 52000},
        ]
        result_path = generator.generate(data_list, output_path, config)

        assert result_path.exists()

    def test_pptx_without_timestamp(
        self, sample_data: pd.DataFrame, output_path: Path
    ) -> None:
        """Test PPTX generation without timestamp."""
        config = PPTXReportConfig(title="No Timestamp", include_timestamp=False)
        generator = PPTXReportGenerator()
        result_path = generator.generate(sample_data, output_path, config)

        assert result_path.exists()

        # Verify presentation has only title + data (no timestamp metadata)
        prs = Presentation(str(result_path))
        assert len(prs.slides) >= 2

    def test_pptx_with_single_column(self, output_path: Path) -> None:
        """Test PPTX with single column DataFrame."""
        single_col_data = pd.DataFrame({"Values": [10, 20, 30, 40]})

        config = PPTXReportConfig(title="Single Column Report")
        generator = PPTXReportGenerator()
        result_path = generator.generate(single_col_data, output_path, config)

        assert result_path.exists()

    def test_pptx_max_rows_per_slide_boundary(self, output_path: Path) -> None:
        """Test pagination boundary conditions."""
        # Exactly max_rows_per_slide (should be 1 data slide)
        exact_data = pd.DataFrame({"Item": [f"Item {i}" for i in range(10)]})

        config = PPTXReportConfig(title="Exact Rows", max_rows_per_slide=10)
        generator = PPTXReportGenerator()
        result_path = generator.generate(exact_data, output_path, config)

        assert result_path.exists()

        # Verify: Title + 1 data slide
        prs = Presentation(str(result_path))
        assert len(prs.slides) == 2

    def test_pptx_chart_with_insufficient_columns(self, output_path: Path) -> None:
        """Test chart generation with insufficient columns (should skip chart)."""
        single_col_data = pd.DataFrame({"Values": [10, 20, 30]})

        config = PPTXReportConfig(
            title="Insufficient Columns", include_chart=True, chart_type="bar"
        )
        generator = PPTXReportGenerator()
        result_path = generator.generate(single_col_data, output_path, config)

        assert result_path.exists()

        # Verify: Title + data only (no chart due to insufficient columns)
        prs = Presentation(str(result_path))
        assert len(prs.slides) == 2  # Title + data (chart skipped)
